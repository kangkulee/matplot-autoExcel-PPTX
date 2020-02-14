from matplotlib import pyplot as plt, font_manager, rc
font_name = font_manager.FontProperties(fname="c:/Windows/Fonts/malgun.ttf").get_name()
rc('font', family=font_name)
plt.rcParams['axes.unicode_minus'] = False
import json
import openpyxl
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from datetime import date

def make_rssi():
    #
    datenow = str(date.today())[2:4] + str(date.today())[5:7] + str(date.today())[8:10]
    name = 'tester'
    filename = 'createCSV/' + datenow + 'test.csv'
    excelname = 'Excel/' + datenow + 'test' + name + '.xlsx'
    picturename = 'chartPicture/avgSnrChart_'
    pptname = 'ppt/' + datenow + 'test' + name + '.pptx'
    titletext = 'title'

    f = open(filename, 'r')

    station_key = dict()

    station_col = 2  # 초기 기지국 별 셀 위치
    xbar_row = 2  # 초기 x축 셀 위치
    cnt_col = 2  # 초기 개수 셀 위치
    cnt_row = 2  # 초기 개수 셀 위치

    count = [] # 개수를 받는 빈 리스트 선언
    #

    # load pptx
    prs = Presentation(pptname)

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "title"
    subtitle.text = titletext

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide_title_shape = slide.shapes.title
    slide_title_shape.text = "STATION LIST("+titletext+")"
    slide_body_shape = slide.placeholders[1]  # text box 객체
    slide_body_tf = slide_body_shape.text_frame

    # load excel
    write_wb = openpyxl.load_workbook(excelname)
    write_ws = write_wb.create_sheet('second')

    for line in f:
        rawdata = json.loads(line.replace("'", "\""))
        if rawdata['payload']['station'] in station_key:
            if rawdata['payload']['rssi'] != 'null':
                station_key[rawdata['payload']['station']].append(float(rawdata['payload']['rssi']))
        else:
            if rawdata['payload']['rssi'] != 'null':
                station_key[rawdata['payload']['station']] = [float(rawdata['payload']['rssi'])]

    write_ws.cell(row=1, column=1, value='x축범위')

    for items in station_key.items():
        avg = round(sum(items[1]) / len(items[1]), 2)

        ys, xs, patches = plt.hist(items[1], range=(-160, 0), color='r', edgecolor='black', linewidth=1.2, rwidth=0.8, bins=160, alpha=0.4)

        count.append(ys)
        write_ws.cell(row=1, column=station_col, value='rssi(station : ' + items[0] + ')')
        station_col = station_col + 1

        plt.xlabel('rssi')
        plt.ylabel('개수')
        plt.title('rssi 차트\n(station : ' + items[0] + ') 평균 : ' + str(avg))

        plt.grid()
        plt.axvline(avg, linestyle='--')

        plt.minorticks_on()

        check_picture = os.path.exists(picturename + items[0] + '.png')
        if check_picture == True:
            # delete picture
            os.unlink(picturename + items[0] + '.png')
            # save picture
            plt.savefig(picturename + items[0] + '.png', dpi=500)
        else:
            # save picture
            plt.savefig(picturename + items[0] + '.png', dpi=500)

        # show chart
        # plt.show()

        plt.clf() # initialize

        # save pptx
        slide_p = slide_body_tf.add_paragraph()
        slide_p.text = items[0]
        slide_p.font.size = Pt(17)
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(picturename + items[0] + '.png', Inches(0.5), Inches(0.5), width=Inches(9), height=Inches(6))

    prs.save(pptname)

    # insert excel
    for i in range(0, len(xs)-1):
        write_ws.cell(row=xbar_row, column=1, value=xs[i])
        xbar_row = xbar_row + 1

    for i in range(len(count)):
        for j in range(len(count[i])):
            write_ws.cell(row=cnt_row, column=cnt_col, value=int(count[i][j]))
            cnt_row = cnt_row + 1
        if cnt_row == 162:
            cnt_row = 2
            cnt_col = cnt_col + 1
    # save excel
    write_wb.save(excelname)
